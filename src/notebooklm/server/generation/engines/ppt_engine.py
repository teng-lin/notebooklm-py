from __future__ import annotations

import json
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..base import ContentGenerator, GeneratedContent, PreviewResult, TemplateInfo


@dataclass
class PptGenerator(ContentGenerator):
    content_type: str = "ppt"
    notebooklm_client: Any = None
    media_root: Path = field(default_factory=lambda: Path.home() / ".notebooklm" / "data")
    _templates: list[dict] = field(default_factory=list)

    def __post_init__(self) -> None:
        self._load_templates()

    def _load_templates(self) -> None:
        template_dir = Path(__file__).resolve().parent.parent / "templates" / "ppt"
        if template_dir.is_dir():
            for f in sorted(template_dir.glob("*.json")):
                try:
                    self._templates.append(json.loads(f.read_text(encoding="utf-8")))
                except Exception:
                    pass

    def _get_template(self, name: str) -> dict:
        for t in self._templates:
            if t["name"] == name:
                return dict(t)
        return dict(self._templates[0]) if self._templates else {
            "name": "default",
            "slide_width": 13.333,
            "slide_height": 7.5,
            "background_color": "#FFFFFF",
            "font_family": "Microsoft YaHei",
            "title_color": "#333333",
            "body_color": "#666666",
            "accent_color": "#1E90FF",
            "layouts": [],
        }

    def _build_slides(self, prs: Presentation, template: dict, slides_data: list[dict]) -> None:
        bg_color = template.get("background_color", "#FFFFFF")
        title_color = template.get("title_color", "#333333")
        body_color = template.get("body_color", "#666666")
        font_family = template.get("font_family", "Microsoft YaHei")

        def _parse_hex(color_str: str) -> RGBColor:
            h = color_str.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        bg_rgb = _parse_hex(bg_color)
        title_rgb = _parse_hex(title_color)
        body_rgb = _parse_hex(body_color)

        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = bg_rgb

            slide_title = slide_data.get("title", "")
            content_items = slide_data.get("items", [])
            layout_type = slide_data.get("layout", "content")

            if layout_type == "title_slide":
                txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(36)
                p.font.color.rgb = title_rgb
                p.font.name = font_family
                p.alignment = PP_ALIGN.CENTER
                if content_items:
                    txBox2 = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = str(content_items[0]) if content_items else ""
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = body_rgb
                    p2.font.name = font_family
                    p2.alignment = PP_ALIGN.CENTER
            elif layout_type == "section_header":
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(28)
                p.font.color.rgb = title_rgb
                p.font.name = font_family
                p.alignment = PP_ALIGN.LEFT
            else:
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(24)
                p.font.color.rgb = title_rgb
                p.font.name = font_family
                p.alignment = PP_ALIGN.LEFT
                txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for i, item in enumerate(content_items):
                    p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                    p2.text = f"\u2022 {item}"
                    p2.font.size = Pt(14)
                    p2.font.color.rgb = body_rgb
                    p2.font.name = font_family
                    p2.space_after = Pt(6)

    async def generate(
        self,
        notebook_id: str,
        prompt: str,
        template: str | None = None,
        options: dict | None = None,
    ) -> GeneratedContent:
        opts = options or {}
        title = opts.get("title", f"PPT {uuid.uuid4().hex[:8]}")
        user_id = opts.get("user_id", 0)
        template_name = template or "classic"
        tmpl = self._get_template(template_name)

        slides_data = opts.get("slides", [])
        if not slides_data:
            slides_data = [
                {"title": title, "items": [prompt], "layout": "title_slide"},
                {"title": "\u5185\u5bb9\u6982\u8ff0", "items": ["\u8981\u70b9 1", "\u8981\u70b9 2", "\u8981\u70b9 3"], "layout": "section_header"},
                {"title": "\u8be6\u7ec6\u5185\u5bb9", "items": ["\u8be6\u7ec6\u8bf4\u660e\u7b2c\u4e00\u70b9", "\u8be6\u7ec6\u8bf4\u660e\u7b2c\u4e8c\u70b9", "\u8be6\u7ec6\u8bf4\u660e\u7b2c\u4e09\u70b9"], "layout": "content"},
                {"title": "\u603b\u7ed3", "items": ["\u5173\u952e\u7ed3\u8bba", "\u4e0b\u4e00\u6b65\u884c\u52a8"], "layout": "content"},
            ]

        prs = Presentation()
        prs.slide_width = Inches(tmpl.get("slide_width", 13.333))
        prs.slide_height = Inches(tmpl.get("slide_height", 7.5))

        self._build_slides(prs, tmpl, slides_data)

        user_dir = self.media_root / "generated" / str(user_id) / notebook_id
        user_dir.mkdir(parents=True, exist_ok=True)

        pptx_filename = f"ppt_{uuid.uuid4().hex[:8]}.pptx"
        pptx_path = user_dir / pptx_filename
        prs.save(str(pptx_path))

        file_size = pptx_path.stat().st_size
        slide_count = len(prs.slides)

        return GeneratedContent(
            id=0,
            content_type="ppt",
            title=title,
            status="completed",
            local_file_path=str(pptx_path),
            file_size=file_size,
            content=json.dumps(slides_data, ensure_ascii=False),
            metadata={
                "ppt_page_count": slide_count,
                "ppt_template": template_name,
                "ppt_json": json.dumps(slides_data, ensure_ascii=False),
            },
        )

    async def preview(
        self,
        notebook_id: str,
        prompt: str,
        template: str | None = None,
    ) -> PreviewResult:
        return PreviewResult(
            content_type="ppt",
            outline=[
                {"title": "\u5c01\u9762", "type": "title_slide"},
                {"title": "\u5185\u5bb9\u5927\u7eb2", "type": "section_header"},
                {"title": "\u8be6\u7ec6\u5185\u5bb9", "type": "content"},
                {"title": "\u603b\u7ed3", "type": "content"},
            ],
            estimated_pages=4,
        )

    async def get_supported_templates(self) -> list[TemplateInfo]:
        if not self._templates:
            return [
                TemplateInfo(name="classic", label="\u7ecf\u5178", description="\u7ecf\u5178\u767d\u5e95\u9ed1\u5b57\u7b80\u7ea6\u6a21\u677f"),
                TemplateInfo(name="modern", label="\u73b0\u4ee3", description="\u84dd\u767d\u914d\u8272\u73b0\u4ee3\u98ce\u683c\u6a21\u677f"),
            ]
        return [
            TemplateInfo(
                name=t["name"],
                label=t.get("label", t["name"]),
                description=t.get("description", ""),
            )
            for t in self._templates
        ]
